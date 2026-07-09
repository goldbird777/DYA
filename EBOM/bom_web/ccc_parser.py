"""
CCC (Color Combination Chart) 파일 파서
PPTX 또는 XLSX에서 CCC 코드 목록을 추출한다.
"""
import re
import openpyxl

# CCC 코드 패턴: 대문자2개+숫자1개+영문숫자1개 (예: A7A, BQ6, C12)
CCC_PATTERN = re.compile(r'^[A-Z]{1,2}\d{1,2}[A-Z0-9]$')

MATERIAL_KEYWORDS = {
    'CLOTH': 'CLOTH',
    'A/CLOTH': 'A/CLOTH',
    'SEMI': 'A/CLOTH',
    'A.LEATHER': 'A.LEATHER',
    'A/LEATHER': 'A.LEATHER',
    'PURE LEATHER': 'PURE LEATHER',
    'PURE': 'PURE LEATHER',
    'LEATHER': 'PURE LEATHER',
}

COLOR_KEYWORDS = {
    'SATURN BLACK': 'SATURN BLACK',
    'BLACK': 'BLACK',
    'NAVY': 'NAVY',
    'GENTLE BROWN': 'GENTLE BROWN',
    'MIDNIGHT GREEN': 'MIDNIGHT GREEN',
    'CARAMEL': 'CARAMEL',
    'GREEN': 'GREEN',
    'BROWN': 'BROWN',
}

MARKET_PATTERNS = re.compile(r'K\d')


def is_ccc_code(val: str) -> bool:
    if not val:
        return False
    v = str(val).strip()
    return bool(CCC_PATTERN.match(v))


def parse_excel(path: str) -> list:
    """
    영업이 관리하는 XLSX에서 CCC 정보 추출.
    행을 스캔하여 CCC 코드, 재질, 시장코드를 파싱.
    """
    wb = openpyxl.load_workbook(path, data_only=True)
    items = []
    seen = set()

    for ws in wb.worksheets:
        current_material = ''
        current_color = ''
        market_cols = {}   # col_idx -> market_code

        for ri, row in enumerate(ws.iter_rows(values_only=True)):
            row_text = ' '.join(str(v) for v in row if v)

            # 시장 코드 헤더 행 감지 (K1, K2, K4 등이 여러 열에 있는 행)
            markets_in_row = MARKET_PATTERNS.findall(row_text)
            if len(markets_in_row) >= 2:
                for ci, cell in enumerate(row):
                    if cell and MARKET_PATTERNS.match(str(cell).strip()):
                        market_cols[ci] = str(cell).strip()
                continue

            # 재질 키워드 감지
            for kw, mat in MATERIAL_KEYWORDS.items():
                if kw in row_text.upper():
                    current_material = mat
                    break

            # 색상 키워드 감지
            for kw, color in COLOR_KEYWORDS.items():
                if kw in row_text.upper():
                    current_color = color
                    break

            # CCC 코드 수집
            ccc_codes_in_row = []
            market_codes_for_row = []

            for ci, cell in enumerate(row):
                if cell is None:
                    continue
                val = str(cell).strip()
                if is_ccc_code(val) and val not in seen:
                    ccc_codes_in_row.append(val)
                    if ci in market_cols:
                        market_codes_for_row.append(market_cols[ci])

            # 행에서 시장코드 컬럼 기반으로 CCC별 시장 매핑
            if ccc_codes_in_row:
                for code in ccc_codes_in_row:
                    seen.add(code)
                    # 해당 코드가 있는 컬럼들의 시장코드 수집
                    mkts = []
                    for ci, cell in enumerate(row):
                        if cell and str(cell).strip() == code:
                            # 같은 행에서 이 코드 왼쪽/오른쪽 컬럼들의 시장코드
                            pass
                    # market_cols가 있으면 그 컬럼 기준으로 매핑
                    if market_cols:
                        for ci, mkt in market_cols.items():
                            if ci < len(row) and row[ci] and is_ccc_code(str(row[ci]).strip()):
                                if str(row[ci]).strip() == code:
                                    mkts.append(mkt)

                    items.append({
                        'ccc_code':      code,
                        'material_type': current_material,
                        'color_name':    current_color,
                        'key_color':     '',
                        'door_code':     '',
                        'stitch_code':   '',
                        'market_codes':  ','.join(mkts) if mkts else '',
                        'remarks':       '',
                    })

    # 중복 제거 후 반환 (CCC 코드 기준)
    unique = {}
    for it in items:
        code = it['ccc_code']
        if code not in unique:
            unique[code] = it
        else:
            # 시장코드 병합
            existing = unique[code]['market_codes']
            new_mkts = it['market_codes']
            merged = set(existing.split(',')) | set(new_mkts.split(','))
            merged.discard('')
            unique[code]['market_codes'] = ','.join(sorted(merged))

    return list(unique.values())


def parse_pptx(path: str) -> list:
    """
    CCC 배색도 PPTX에서 CCC 코드 및 재질/색상/시장코드 추출.
    슬라이드 내 테이블을 스캔.
    """
    try:
        from pptx import Presentation
    except ImportError:
        return []

    prs = Presentation(path)
    items = []
    seen = set()

    for slide in prs.slides:
        # 슬라이드 제목에서 재질 파악
        slide_material = ''
        for shape in slide.shapes:
            if shape.has_text_frame:
                txt = shape.text_frame.text.upper()
                for kw, mat in MATERIAL_KEYWORDS.items():
                    if kw in txt:
                        slide_material = mat
                        break

        for shape in slide.shapes:
            if not shape.has_table:
                continue
            tbl = shape.table
            if tbl.rows is None or len(tbl.rows) < 2:
                continue

            # 헤더 행에서 컬럼 역할 파악
            header = [c.text.strip() for c in tbl.rows[0].cells]
            # SEAT CODE 컬럼 찾기
            seat_col = next((i for i, h in enumerate(header) if 'SEAT CODE' in h.upper() or 'CODE' in h.upper()), None)
            name_col = next((i for i, h in enumerate(header) if 'NAME' in h.upper()), None)
            door_col = next((i for i, h in enumerate(header) if 'DOOR' in h.upper()), None)
            stitch_col = next((i for i, h in enumerate(header) if 'STITCH' in h.upper()), None)
            remarks_col = next((i for i, h in enumerate(header) if 'REMARK' in h.upper()), None)

            if seat_col is None:
                continue

            current_color = ''
            current_key = ''

            for row in list(tbl.rows)[1:]:
                cells = [c.text.strip().replace('\n', ' ') for c in row.cells]
                if len(cells) <= seat_col:
                    continue

                # 색상 이름, KEY COLOR 수집
                if name_col is not None and cells[name_col]:
                    txt = cells[name_col].upper()
                    for kw, color in COLOR_KEYWORDS.items():
                        if kw in txt:
                            current_color = color
                            break
                # 두 번째 열에 KEY COLOR CODE 있는 경우 (NAME | CODE | SEAT CODE ...)
                if name_col is not None and name_col + 1 < len(cells):
                    key_candidate = cells[name_col + 1].strip()
                    if re.match(r'^[A-Z0-9]{2,4}$', key_candidate):
                        current_key = key_candidate

                ccc_val = cells[seat_col].strip()
                if not is_ccc_code(ccc_val):
                    continue
                if ccc_val in seen:
                    continue
                seen.add(ccc_val)

                door = cells[door_col].strip() if door_col is not None and door_col < len(cells) else ''
                stitch = cells[stitch_col].strip() if stitch_col is not None and stitch_col < len(cells) else ''
                remarks = cells[remarks_col].strip() if remarks_col is not None and remarks_col < len(cells) else ''

                items.append({
                    'ccc_code':      ccc_val,
                    'material_type': slide_material,
                    'color_name':    current_color,
                    'key_color':     current_key,
                    'door_code':     door,
                    'stitch_code':   stitch,
                    'market_codes':  '',
                    'remarks':       remarks,
                })

    return items


def parse_ccc_file(path: str, ext: str) -> list:
    ext = ext.lower().lstrip('.')
    if ext in ('xlsx', 'xlsm'):
        return parse_excel(path)
    elif ext == 'pptx':
        return parse_pptx(path)
    return []
